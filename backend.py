"""
AI Study Companion - Backend Module
Lightweight, deterministic implementations that work offline
"""

import re
import json
import sqlite3
from pathlib import Path
from datetime import datetime, timedelta
from collections import Counter
import pandas as pd
import numpy as np

import nltk

# Auto-download required NLTK data if missing
for resource in ["punkt", "punkt_tab"]:
    try:
        nltk.data.find(f"tokenizers/{resource}")
    except LookupError:
        nltk.download(resource)

# Try to import heavy libraries, but provide fallbacks
try:
    from transformers import pipeline
    TRANSFORMERS_AVAILABLE = True
except ImportError:
    TRANSFORMERS_AVAILABLE = False

try:
    import nltk
    from nltk.corpus import stopwords
    from nltk.tokenize import sent_tokenize, word_tokenize
    NLTK_AVAILABLE = True
except ImportError:
    NLTK_AVAILABLE = False
    # Provide basic fallback tokenization
    def sent_tokenize(text):
        return re.split(r'[.!?]+', text)
    def word_tokenize(text):
        return text.split()

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    SKLEARN_AVAILABLE = True
except ImportError:
    SKLEARN_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from PyPDF2 import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class FileExtractor:
    """Extract text from various file formats with robust error handling."""
    
    @staticmethod
    def extract_text(file_path):
        """
        Extract text from file with fallbacks.
        
        Args:
            file_path: Path to file
            
        Returns:
            str: Extracted text or error message
        """
        try:
            file_path = Path(file_path)
            extension = file_path.suffix.lower()
            
            if extension == '.txt':
                return FileExtractor._extract_txt(file_path)
            elif extension in ['.md', '.markdown']:
                return FileExtractor._extract_txt(file_path)
            elif extension == '.pdf':
                return FileExtractor._extract_pdf(file_path)
            elif extension in ['.docx', '.doc']:
                return FileExtractor._extract_docx(file_path)
            else:
                return f"Unsupported file type: {extension}"
        except Exception as e:
            return f"Error extracting text: {str(e)}"
    
    @staticmethod
    def _extract_txt(file_path):
        """Extract from text file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            raise Exception(f"Text extraction failed: {e}")
    
    @staticmethod
    def _extract_pdf(file_path):
        """Extract from PDF with multiple fallbacks."""
        text = ""
        
        # Try pdfplumber first (best quality)
        if PDFPLUMBER_AVAILABLE:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                if text.strip():
                    return text
            except Exception as e:
                pass
        
        # Fallback to PyPDF2
        if PYPDF2_AVAILABLE:
            try:
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                if text.strip():
                    return text
            except Exception as e:
                pass
        
        # If no PDF library available or extraction failed
        return "PDF extraction not available. Please install pdfplumber or PyPDF2, or use a text file."
    
    @staticmethod
    def _extract_docx(file_path):
        """Extract from DOCX."""
        if not DOCX_AVAILABLE:
            return "DOCX extraction not available. Please install python-docx or use a text file."
        
        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {e}")
    
    @staticmethod
    def clean_text(text):
        """Clean and normalize text."""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters (keep basic punctuation)
        text = re.sub(r'[^a-zA-Z0-9\s.,;:!?\'"()-]', '', text)
        return text.strip()
    
    @staticmethod
    def get_preview(text, max_words=200):
        """Get preview of text."""
        words = text.split()
        if len(words) <= max_words:
            return text
        return ' '.join(words[:max_words]) + '...'
    
    @staticmethod
    def get_text_stats(text):
        """Get text statistics."""
        if not text or not text.strip():
            return {
                'characters': 0,
                'words': 0,
                'sentences': 0,
                'avg_word_length': 0,
                'avg_sentence_length': 0
            }
        
        words = text.split()
        sentences = sent_tokenize(text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        return {
            'characters': len(text),
            'words': len(words),
            'sentences': len(sentences),
            'avg_word_length': sum(len(w) for w in words) / max(len(words), 1),
            'avg_sentence_length': len(words) / max(len(sentences), 1)
        }
    
    @staticmethod
    def chunk_text(text, chunk_size=1000, overlap=100):
        """Split text into overlapping chunks."""
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if chunk:
                chunks.append(chunk)
        
        return chunks


class SmartSummarizer:
    """Hybrid summarizer with extractive and abstractive modes."""
    
    def __init__(self, device='cpu'):
        self.device = device
        self.cache = {}
        self.abstractive_model = None
        
        # Load stopwords if NLTK available
        if NLTK_AVAILABLE:
            try:
                self.stop_words = set(stopwords.words('english'))
            except:
                nltk.download('stopwords', quiet=True)
                self.stop_words = set(stopwords.words('english'))
        else:
            # Basic English stopwords
            self.stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
                             'of', 'with', 'by', 'from', 'as', 'is', 'was', 'are', 'were', 'be',
                             'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will',
                             'would', 'could', 'should', 'may', 'might', 'can'}
    
    def summarize(self, text, method='extractive', num_sentences=5, max_length=150):
        """
        Summarize text.
        
        Args:
            text: Input text
            method: 'extractive' or 'abstractive'
            num_sentences: Number of sentences for extractive
            max_length: Max length for abstractive
            
        Returns:
            str: Summary
        """
        if not text or len(text.strip()) < 50:
            return "Text too short to summarize."
        
        # Check cache
        cache_key = f"{method}_{num_sentences}_{max_length}_{hash(text)}"
        if cache_key in self.cache:
            return self.cache[cache_key]
        
        try:
            if method == 'extractive':
                summary = self._extractive_summary(text, num_sentences)
            elif method == 'abstractive':
                summary = self._abstractive_summary(text, max_length)
            else:
                summary = self._extractive_summary(text, num_sentences)
            
            # Cache result
            self.cache[cache_key] = summary
            return summary
            
        except Exception as e:
            return f"Summarization error: {str(e)}. Using first {num_sentences} sentences as fallback."
    
    def _extractive_summary(self, text, num_sentences):
        """Fast extractive summarization using sentence scoring."""
        # Tokenize sentences
        sentences = sent_tokenize(text)
        if len(sentences) <= num_sentences:
            return text
        
        # Tokenize words
        words = word_tokenize(text.lower())
        
        # Filter stopwords
        filtered_words = [w for w in words if w.isalnum() and w not in self.stop_words]
        
        # Calculate word frequencies
        word_freq = Counter(filtered_words)
        
        # Score sentences
        sentence_scores = {}
        for i, sentence in enumerate(sentences):
            sentence_words = word_tokenize(sentence.lower())
            score = sum(word_freq.get(w, 0) for w in sentence_words if w in word_freq)
            # Normalize by sentence length
            if len(sentence_words) > 0:
                sentence_scores[i] = score / len(sentence_words)
        
        # Get top sentences
        top_indices = sorted(sentence_scores, key=sentence_scores.get, reverse=True)[:num_sentences]
        
        # Return in original order
        top_indices.sort()
        summary = ' '.join(sentences[i] for i in top_indices)
        
        return summary
    
    def _abstractive_summary(self, text, max_length):
        """Abstractive summarization using transformer model."""
        if not TRANSFORMERS_AVAILABLE:
            return "Abstractive summarization requires transformers library. Using extractive instead:\n\n" + \
                   self._extractive_summary(text, 5)
        
        # Load model lazily
        if self.abstractive_model is None:
            try:
                self.abstractive_model = pipeline(
                    "summarization",
                    model="facebook/bart-large-cnn",
                    device=0 if self.device == 'cuda' else -1
                )
            except Exception as e:
                return f"Model loading failed: {e}. Using extractive instead:\n\n" + \
                       self._extractive_summary(text, 5)
        
        try:
            # Chunk if too long
            max_input_length = 1024
            if len(text.split()) > max_input_length:
                chunks = FileExtractor.chunk_text(text, max_input_length, 50)
                summaries = []
                for chunk in chunks[:3]:  # Limit chunks
                    result = self.abstractive_model(
                        chunk,
                        max_length=max_length,
                        min_length=30,
                        do_sample=False
                    )
                    summaries.append(result[0]['summary_text'])
                return ' '.join(summaries)
            else:
                result = self.abstractive_model(
                    text,
                    max_length=max_length,
                    min_length=30,
                    do_sample=False
                )
                return result[0]['summary_text']
        except Exception as e:
            return f"Abstractive summarization failed: {e}. Using extractive:\n\n" + \
                   self._extractive_summary(text, 5)


class TopicExtractor:
    """Extract key topics and keywords from text."""
    
    def __init__(self):
        if NLTK_AVAILABLE:
            try:
                self.stop_words = set(stopwords.words('english'))
            except:
                nltk.download('stopwords', quiet=True)
                self.stop_words = set(stopwords.words('english'))
        else:
            self.stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for'}
    
    def extract_keywords(self, text, top_n=10):
        """Extract keywords using TF-IDF or frequency."""
        if SKLEARN_AVAILABLE:
            return self._tfidf_keywords(text, top_n)
        else:
            return self._simple_keywords(text, top_n)
    
    def _tfidf_keywords(self, text, top_n):
        """TF-IDF based keyword extraction."""
        try:
            vectorizer = TfidfVectorizer(
                max_features=top_n * 2,
                stop_words='english',
                ngram_range=(1, 2)
            )
            tfidf_matrix = vectorizer.fit_transform([text])
            feature_names = vectorizer.get_feature_names_out()
            scores = tfidf_matrix.toarray()[0]
            
            keywords = [(feature_names[i], scores[i]) for i in range(len(feature_names))]
            keywords.sort(key=lambda x: x[1], reverse=True)
            
            return keywords[:top_n]
        except:
            return self._simple_keywords(text, top_n)
    
    def _simple_keywords(self, text, top_n):
        """Simple frequency-based keywords."""
        words = word_tokenize(text.lower())
        filtered = [w for w in words if w.isalnum() and len(w) > 3 and w not in self.stop_words]
        freq = Counter(filtered)
        return freq.most_common(top_n)
    
    def extract_topics(self, text, num_topics=5):
        """Extract main topics."""
        keywords = self.extract_keywords(text, top_n=num_topics)
        return [kw[0].title() for kw in keywords]
    
    def generate_tags(self, text, max_tags=7):
        """Generate hashtag-style tags."""
        topics = self.extract_topics(text, num_topics=max_tags)
        return [f"#{topic.replace(' ', '')}" for topic in topics]


class SlideGenerator:
    """Generate presentation slides from content."""
    
    def generate_slide_content(self, text, summary, keywords, num_slides=5):
        """Generate structured slide content."""
        slides = []
        
        # Title slide
        slides.append({
            'type': 'title',
            'title': keywords[0] if keywords else 'Study Notes',
            'subtitle': 'Key Concepts and Summary'
        })
        
        # Keywords slide
        slides.append({
            'type': 'bullet',
            'title': 'Key Topics',
            'bullets': keywords[:6]
        })
        
        # Summary slides
        sentences = sent_tokenize(summary)
        chunks = [sentences[i:i+3] for i in range(0, len(sentences), 3)]
        
        for i, chunk in enumerate(chunks[:num_slides-2]):
            slides.append({
                'type': 'bullet',
                'title': f'Summary Part {i+1}',
                'bullets': chunk
            })
        
        return slides
    
    def create_presentation(self, slide_content, output_path=None):
        """Create PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            # Create a simple text file instead
            if output_path is None:
                output_path = Path('presentation.txt')
            
            with open(output_path, 'w') as f:
                f.write("PRESENTATION SLIDES\n")
                f.write("=" * 50 + "\n\n")
                for i, slide in enumerate(slide_content, 1):
                    f.write(f"Slide {i}: {slide['title']}\n")
                    f.write("-" * 50 + "\n")
                    if 'subtitle' in slide:
                        f.write(f"{slide['subtitle']}\n")
                    if 'bullets' in slide:
                        for bullet in slide['bullets']:
                            f.write(f"  • {bullet}\n")
                    f.write("\n")
            
            return str(output_path)
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for slide_data in slide_content:
                if slide_data['type'] == 'title':
                    self._add_title_slide(prs, slide_data)
                elif slide_data['type'] == 'bullet':
                    self._add_bullet_slide(prs, slide_data)
            
            if output_path is None:
                output_path = Path('presentation.pptx')
            
            prs.save(output_path)
            return str(output_path)
            
        except Exception as e:
            raise Exception(f"Presentation creation failed: {e}")
    
    def _add_title_slide(self, prs, slide_data):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = slide_data['title']
        subtitle.text = slide_data.get('subtitle', '')
        
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
    
    def _add_bullet_slide(self, prs, slide_data):
        """Add bullet slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        text_frame = body.text_frame
        text_frame.clear()
        
        for bullet in slide_data['bullets']:
            p = text_frame.add_paragraph()
            p.text = str(bullet)
            p.level = 0
            p.font.size = Pt(20)


class QuizGenerator:
    """Generate practice quizzes from text."""
    
    def __init__(self):
        self.question_templates = [
            "What is {}?",
            "Which of the following best describes {}?",
            "What are the main characteristics of {}?",
            "How does {} work?",
            "What is the purpose of {}?"
        ]
    
    def generate_quiz(self, text, num_questions=5, difficulty='medium'):
        """Generate quiz questions."""
        try:
            # Extract keywords as concepts
            extractor = TopicExtractor()
            keywords = extractor.extract_keywords(text, top_n=num_questions * 2)
            sentences = sent_tokenize(text)
            
            questions = []
            used_concepts = set()
            
            for keyword, score in keywords:
                if len(questions) >= num_questions:
                    break
                
                if keyword in used_concepts:
                    continue
                
                # Find relevant sentence
                relevant_sentences = [s for s in sentences if keyword.lower() in s.lower()]
                
                if not relevant_sentences:
                    continue
                
                question = self._create_question(keyword, relevant_sentences[0], difficulty)
                if question:
                    questions.append(question)
                    used_concepts.add(keyword)
            
            # Fill remaining with definition questions
            while len(questions) < num_questions and len(used_concepts) < len(keywords):
                for keyword, score in keywords:
                    if keyword not in used_concepts and len(questions) < num_questions:
                        question = self._create_simple_question(keyword, text)
                        if question:
                            questions.append(question)
                            used_concepts.add(keyword)
            
            return questions
            
        except Exception as e:
            # Return sample questions as fallback
            return self._generate_fallback_quiz(num_questions)
    
    def _create_question(self, concept, context, difficulty):
        """Create a quiz question."""
        import random
        
        template = random.choice(self.question_templates)
        question_text = template.format(concept.title())
        
        # Simplify context for answer
        answer = context[:100] + "..." if len(context) > 100 else context
        
        # Generate distractors
        distractors = [
            f"A method for processing {concept}",
            f"An approach to {concept} implementation",
            f"A framework for {concept} management"
        ]
        
        options = [answer] + distractors[:2]
        random.shuffle(options)
        
        return {
            'question': question_text,
            'options': options,
            'correct_answer': options.index(answer),
            'explanation': context,
            'concept': concept,
            'difficulty': difficulty
        }
    
    def _create_simple_question(self, concept, text):
        """Create simple definition question."""
        import random
        
        question_text = f"What is {concept.title()}?"
        
        answer = f"Related to {concept}"
        distractors = [
            "A programming language",
            "A data structure",
            "An algorithm technique"
        ]
        
        options = [answer] + distractors[:2]
        random.shuffle(options)
        
        return {
            'question': question_text,
            'options': options,
            'correct_answer': options.index(answer),
            'explanation': f"The text discusses {concept}",
            'concept': concept,
            'difficulty': 'easy'
        }
    
    def _generate_fallback_quiz(self, num_questions):
        """Generate sample quiz as fallback."""
        sample_questions = [
            {
                'question': 'What is the main topic of this text?',
                'options': ['General concepts', 'Specific details', 'Historical context', 'Future predictions'],
                'correct_answer': 0,
                'explanation': 'Based on the content provided',
                'concept': 'general',
                'difficulty': 'easy'
            },
            {
                'question': 'Which key concept is discussed?',
                'options': ['Primary concept', 'Secondary concept', 'Tertiary concept', 'None of these'],
                'correct_answer': 0,
                'explanation': 'The text focuses on primary concepts',
                'concept': 'concepts',
                'difficulty': 'medium'
            }
        ]
        
        return sample_questions[:num_questions]
    
    def export_quiz(self, questions, format='json', output_path=None):
        """Export quiz to file."""
        if output_path is None:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            ext = 'json' if format == 'json' else 'txt'
            output_path = f'quiz_{timestamp}.{ext}'
        
        if format == 'json':
            with open(output_path, 'w') as f:
                json.dump(questions, f, indent=2)
        else:
            with open(output_path, 'w') as f:
                for i, q in enumerate(questions, 1):
                    f.write(f"Question {i}: {q['question']}\n")
                    for j, opt in enumerate(q['options'], 1):
                        f.write(f"  {j}. {opt}\n")
                    f.write(f"  Correct: {q['correct_answer'] + 1}\n\n")
        
        return str(output_path)


class StudyScheduler:
    """Generate optimized study schedules."""
    
    def generate_schedule(self, goals, deadline, daily_hours=4, session_length=45):
        """Generate study schedule."""
        try:
            today = datetime.now().date()
            
            if isinstance(deadline, str):
                deadline_date = datetime.strptime(deadline, '%Y-%m-%d').date()
            else:
                deadline_date = deadline
            
            days_available = max((deadline_date - today).days, 1)
            
            # Estimate time per goal (in minutes)
            time_per_goal = 60  # Default 1 hour per goal
            total_minutes = len(goals) * time_per_goal
            
            # Calculate sessions per day
            sessions_per_day = min(
                len(goals) / days_available,
                (daily_hours * 60) / session_length
            )
            sessions_per_day = max(1, int(sessions_per_day))
            
            # Generate schedule
            schedule = []
            current_date = today
            goal_index = 0
            
            while goal_index < len(goals) and current_date <= deadline_date:
                # Skip weekends (optional)
                if current_date.weekday() < 5:  # Monday-Friday
                    for session_num in range(sessions_per_day):
                        if goal_index >= len(goals):
                            break
                        
                        start_hour = 9 + (session_num * 2)  # Space sessions
                        if start_hour >= 22:
                            break
                        
                        schedule.append({
                            'date': current_date.strftime('%Y-%m-%d'),
                            'start_time': f"{current_date.strftime('%Y-%m-%d')} {start_hour:02d}:00",
                            'end_time': f"{current_date.strftime('%Y-%m-%d')} {start_hour:02d}:{session_length:02d}",
                            'duration': session_length,
                            'task': goals[goal_index],
                            'session_type': 'study',
                            'priority': self._calculate_priority(deadline_date, current_date)
                        })
                        
                        goal_index += 1
                
                current_date += timedelta(days=1)
            
            return schedule
            
        except Exception as e:
            # Return sample schedule
            return [{
                'date': datetime.now().strftime('%Y-%m-%d'),
                'start_time': f"{datetime.now().strftime('%Y-%m-%d')} 09:00",
                'end_time': f"{datetime.now().strftime('%Y-%m-%d')} 10:00",
                'duration': 60,
                'task': 'Study session',
                'session_type': 'study',
                'priority': 'medium'
            }]
    
    def _calculate_priority(self, deadline, current_date):
        """Calculate task priority."""
        days_until = (deadline - current_date).days
        
        if days_until <= 1:
            return 'urgent'
        elif days_until <= 3:
            return 'high'
        elif days_until <= 7:
            return 'medium'
        else:
            return 'low'
    
    def export_schedule(self, schedule, format='csv', output_path=None):
        """Export schedule."""
        if output_path is None:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            ext = 'csv' if format == 'csv' else 'ics'
            output_path = f'schedule_{timestamp}.{ext}'
        
        if format == 'csv':
            df = pd.DataFrame(schedule)
            df.to_csv(output_path, index=False)
        elif format == 'ics':
            with open(output_path, 'w') as f:
                f.write("BEGIN:VCALENDAR\n")
                f.write("VERSION:2.0\n")
                f.write("PRODID:-//AI Study Companion//EN\n")
                
                for session in schedule:
                    f.write("BEGIN:VEVENT\n")
                    f.write(f"SUMMARY:{session['task']}\n")
                    f.write(f"DTSTART:{session['start_time'].replace('-', '').replace(':', '').replace(' ', 'T')}00\n")
                    f.write(f"DTEND:{session['end_time'].replace('-', '').replace(':', '').replace(' ', 'T')}00\n")
                    f.write("END:VEVENT\n")
                
                f.write("END:VCALENDAR\n")
        
        return str(output_path)


class ProgressTracker:
    """Track learning progress and mastery."""
    
    def __init__(self):
        self.metrics = {
            'tasks_completed': 0,
            'tasks_total': 0,
            'quiz_scores': [],
            'study_time_minutes': 0,
            'streak_days': 0
        }
    
    def update_metrics(self, tasks_completed=0, tasks_total=0, quiz_score=None, study_time=0):
        """Update metrics."""
        if tasks_completed:
            self.metrics['tasks_completed'] += tasks_completed
        if tasks_total:
            self.metrics['tasks_total'] = max(self.metrics['tasks_total'], tasks_total)
        if quiz_score is not None:
            self.metrics['quiz_scores'].append(quiz_score)
        if study_time:
            self.metrics['study_time_minutes'] += study_time
    
    def calculate_mastery(self):
        """Calculate mastery percentage."""
        # Task completion (40%)
        task_score = (self.metrics['tasks_completed'] / max(self.metrics['tasks_total'], 1)) * 40
        
        # Quiz performance (40%)
        if self.metrics['quiz_scores']:
            quiz_avg = sum(self.metrics['quiz_scores']) / len(self.metrics['quiz_scores'])
            quiz_score = quiz_avg * 0.4
        else:
            quiz_score = 0
        
        # Time investment (20%)
        expected_time = self.metrics['tasks_total'] * 60
        time_score = min(self.metrics['study_time_minutes'] / max(expected_time, 1), 1) * 20
        
        total_mastery = task_score + quiz_score + time_score
        
        return {
            'overall_mastery': round(total_mastery, 1),
            'task_completion': round((self.metrics['tasks_completed'] / max(self.metrics['tasks_total'], 1)) * 100, 1),
            'quiz_performance': round(quiz_avg if self.metrics['quiz_scores'] else 0, 1),
            'time_investment': round((self.metrics['study_time_minutes'] / max(expected_time, 1)) * 100, 1),
            'tasks_completed': self.metrics['tasks_completed'],
            'tasks_remaining': self.metrics['tasks_total'] - self.metrics['tasks_completed'],
            'avg_quiz_score': round(sum(self.metrics['quiz_scores']) / len(self.metrics['quiz_scores']), 1) if self.metrics['quiz_scores'] else 0,
            'total_study_hours': round(self.metrics['study_time_minutes'] / 60, 1)
        }
    
    def get_recommendations(self):
        """Get study recommendations."""
        mastery = self.calculate_mastery()
        recommendations = []
        
        if mastery['overall_mastery'] < 50:
            recommendations.append("📚 Focus on completing more tasks")
            recommendations.append("⏰ Increase daily study time")
        
        if mastery.get('quiz_performance', 0) < 70:
            recommendations.append("🎯 Review weak topics")
            recommendations.append("📝 Practice with more quizzes")
        
        if mastery['task_completion'] < 80:
            recommendations.append("✅ Prioritize pending tasks")
        
        if mastery['overall_mastery'] >= 80:
            recommendations.append("🎉 Excellent progress!")
            recommendations.append("🚀 Ready for advanced topics")
        
        return recommendations


# Export all classes
__all__ = [
    'FileExtractor',
    'SmartSummarizer',
    'TopicExtractor',
    'SlideGenerator',
    'QuizGenerator',
    'StudyScheduler',
    'ProgressTracker'
]
